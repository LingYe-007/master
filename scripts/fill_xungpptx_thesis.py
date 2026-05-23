#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""将「勋高-202333248.pptx」模板占位符替换为伍勋高论文答辩用语（尽量保留版式，仅改 text_frame.text）。"""
from __future__ import annotations

import re
import shutil
import sys
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "勋高-202333248.pptx"
DST = ROOT / "勋高-202333248-已填内容.pptx"

L1 = "单击此处添加此章节的简要内容，本模板内所有素材均可自由编辑及移动替换，框架完整，套用性强。"
L2 = "单击此处添加此章节的简要内容，本模板内所有素材均可自由编辑及移动替换，框架完整，套用性强。单击此处添加此章节的简要内容，本模板内所有素材均可自由编辑及移动替换，框架完整，套用性强。"
L3 = (
    "单击此处添加此章节的简要内容，本模板内所有素材均可自由编辑及移动替换，框架完整，套用性强。"
    "单击此处添加此章节的简要内容，本模板内所有素材均可自由编辑及移动替换，框架完整，套用性强。"
    "单击此处添加此章节的简要内容，本模板内所有素材均可自由编辑及移动替换，框架完整，套用性强。"
)
L_LOREM = "添加适当的文字，一页的文字最好不要超过200,添加适当的文字，添加适当的文字"
L_LOREM_SHORT = "添加适当的文字，一页的文字最好不要超过200,添加适当的文字"
L_PARA = "这里输入段落文本内容这里输入段落文本内容这里输入段落文本内容这里输入段落文本内容"
L_COLOR = "模板中的所有颜色板式均可编辑替换，可以根据自己的需要更换自己喜欢的颜色与修改图表的布局。"
L_REF = "（论文参考文献）根据需要添加适当的文字，一页的文字最好不要超过200字根据需要添加适当的文字"
L_CLICK = (
    "点击输入简要文字介绍，文字内容需概括精炼，不用多余的文字修饰，言简意赅的说明分项内容点击输入简要文字介绍，"
    "文字内容需概括精炼，不用多余的文字修点击输入简要文字介绍"
)
L_YOUR = "您的内容打在这里，或者通过复制您的文本后，在此框中选择粘贴"
L_YOUR2 = (
    "您的内容打在这里，或者通过复制您的文本后，在此框中选择粘贴，并选择只保留文字。"
    "您的内容打在这里，或者通过复制您的文本后，在此框中选择粘贴，并选择只保留文字。"
)
L_YOUR1 = (
    "您的内容打在这里，或者通过复制您的文本后，在此框中选择粘贴，并选择只保留文字。"
)
L_ADD_TITLE = "添加标题"
L_BAD_TITLE = "添加本本标题"
L_MISC = "添加适当的文字，一页的文字最好不要超过200,添加适当的文字，一页的文字最好不要超过200"
L_MISC2 = "添加适当的文字，最好不要超过200字,添加适当的文字，最好不要超过200字"
L_THESIS_CONC = (
    "论文结论：单击此处添加此章节的简要内容，本模板内所有素材均可自由编辑及移动替换，框架完整，套用性强。"
    "单击此处添加此章节的简要内容，本模板内所有素材均可自由编辑及移动替换，框架完整，套用性强。"
)
L_LONG_BAD = (
    "添加适当的文字，一页的文字最好不要超过200,添加适当的文字，添加适当的文字，一页的文字最好不要超过200,"
    "添加适当的文字，添加适当的文字，一页的文字最好不要超过200,添加适当的文字添加适当的文字，一页的文字最好不要超过200,"
    "添加适当的文字，添加适当的文字，一页的文字最好不要超过200,添加适当的文字，添加适当的文字，一页的文字最好不要超过200,添加适当的文字"
)
FEAS = "添加适当文字，一页文字不要超200添加适当文字,一页文字不要超200"


def set_shape_text(shape, new: str) -> None:
    if not hasattr(shape, "text_frame"):
        return
    shape.text_frame.text = new


def walk_shapes(shapes, fn) -> None:
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            walk_shapes(sh.shapes, fn)
        elif hasattr(sh, "text_frame"):
            fn(sh)


def main() -> int:
    if not SRC.is_file():
        print("找不到源文件:", SRC, file=sys.stderr)
        return 1
    shutil.copy2(SRC, DST)
    prs = Presentation(str(DST))

    # 每页「纯 L1」出现次数（用于顺序填充）
    plain_l1: dict[int, int] = defaultdict(int)

    def global_sub(t: str) -> str:
        t = t.replace("学论网-专注原创-www.xuelun.me", "")
        t = t.replace("简约风毕业论文答辩PPT", "基于对比学习的联邦异构图推荐算法与系统研究")
        t = t.replace("XXXXXXX大学XXXXX学院", "西北大学信息科学与技术学院")
        t = t.replace("答辩人：小雪PPT", "答辩人：伍勋高")
        t = t.replace("指导老师：小雪PPT", "指导老师：吴昊")
        t = t.replace("小雪PPT", "伍勋高")
        t = t.replace("添加本本标题", "小节要点")
        return t

    # 按页、按占位类型顺序填充（0-based slide index）
    slide_fills: dict[int, dict[str, list[str]]] = {
        3: {
            "01": [
                "01\n合规与动机：GDPR/PIPL 与数据主权下，用户原始数据不宜集中；联邦学习以「数据不动模型动」支撑推荐协同。",
                "02\n三难并存：冷启动（本地交互稀疏）、Non-IID（客户端漂移）、通信与异构图模型落地难。",
                "03\n本文路线：FedASCL 提升效果 → 语义感知压缩（元路径剪枝 + 量化）降通信 → 端云系统验证。",
            ]
        },
        4: {
            "plain_l1": [
                "理论意义：在联邦约束下显式利用属性与异构语义，为稀疏与倾斜数据提供可交换表征学习范式。",
                "方法意义：双视图对比 + 全局原型对齐，为缓解冷启动与漂移提供可复现算法框架。",
                "工程意义：分层系统与压缩上传闭环，使算法结论与实测通信、耗时相互印证。",
                "应用意义：面向学术资源推荐场景，兼顾隐私、精度与带宽受限部署。",
            ]
        },
        5: {
            "double_l1": [
                "图/异构推荐：GNN 与元路径在集中式场景成熟，联邦端属性与结构难协同，零交互下图退化。",
                "冷启动：大模型与文本算力路线与端侧图联邦设定不同；需属性语义迁入交互空间。",
                "Non-IID：客户端漂移与全局收敛困难；需跨客户端语义参照（原型等）。",
            ]
        },
    }

    def fill_slide_4(t: str, si: int) -> str:
        if si != 3:
            return t
        for i, prefix in enumerate(["01\n", "02\n", "03\n"]):
            old = prefix + L1
            if t == old:
                return slide_fills[3]["01"][i]
        return t

    def fill_slide_5_plain(t: str, si: int) -> str:
        if si != 4 or t.strip() != L1:
            return t
        idx = plain_l1[si]
        plain_l1[si] += 1
        opts = slide_fills[4]["plain_l1"]
        return opts[idx] if idx < len(opts) else opts[-1]

    def fill_slide_6_double(t: str, si: int) -> str:
        if si != 5:
            return t
        d = L1 + "\n" + L1
        if t.strip() != d:
            return t
        idx = plain_l1[si]
        plain_l1[si] += 1
        opts = slide_fills[5]["double_l1"]
        return opts[idx] if idx < len(opts) else opts[-1]

    # 创新点（slide 7 index 6）组合内 4 段 L1
    innov = [
        "FedASCL：客户端属性语义超图 + 交互结构图双编码；InfoNCE 跨视图对齐；服务器聚合全局语义原型正则。",
        "语义感知压缩：动态元路径 Top-K 通道剪枝；残差梯度量化 + 本地误差补偿；约 20× 压缩下精度可比。",
        "系统实现：展示—业务—算法引擎—存储分层；模型分发—本地训练—压缩上传—安全聚合闭环。",
        "实验与口径：ML-1M / Yelp / ACM；主表、冷启动表、Non-IID 与消融；表间说明见 tab:chap3_exp_table_guide。",
    ]

    # slide 8 (7) 与 slide 9 (8) — 多种「您的内容」占位
    slide8_bodies = [
        "创新点一：双视图对比学习将属性语义迁入稀疏交互空间，配合 BPR 元路径主任务联合优化。",
        "创新点二：本地原型加噪上传、服务器加权聚合为全局原型，缓解 Non-IID 语义漂移。",
        "创新点三：上传阶段串联元路径选择器与残差量化，控制通信且保护异构语义通道。",
        "创新点四：联邦论文推荐系统集成 FedASCL 与压缩，完成界面与实测验证。",
    ]
    slide9_boxes = [
        "01\n算法：FedASCL 各模块与 alg:fedasclrc 流程；复杂度与第三章实验配置、基线、RQ1–3、消融。",
        "02\n压缩：压缩流程、选择器、量化与补偿；收敛/消融/ρ 扫描/帕累托与通信表。",
        "03\n系统：需求与架构、子图不出域、仿真引擎与冷启动路由、界面与隐私模块。",
        "04\n总结：三条贡献与展望（动态图、TEE 隐私、跨域少样本）；致谢与 Q&A 指引。",
    ]

    # 2.6 实验步骤 5 条
    exp_steps = [
        "数据与划分：三数据集、属性表、Dirichlet 客户端 Non-IID；每轮参与比例与配置表。",
        "基线与实现：FedNCF/GNN/Per/HGNN/Proto + 本文；环境与超参按第三章表。",
        "训练与评测：100 轮曲线；HR/NDCG@10；冷启动与 α 扫描及消融。",
        "压缩实验：更难联邦设定；20× 压缩下主指标与通信、挂钟时间。",
        "系统实测：与第四章通信量级对照；功能与监控界面验证。",
    ]

    # 2.7 研究目标 4 条
    goals = [
        "01. 研究目标\n建立联邦异构图推荐在冷启动与 Non-IID 下仍稳健的表示学习机制。",
        "02. 研究目标\n在保持精度的前提下显著降低端云上传通信量。",
        "03. 研究目标\n形成可演示、可复用的端云联邦推荐系统实现。",
        "04. 研究目标\n在统一口径下完成主实验、消融与表间对照说明。",
    ]

    # 关键技术一~四（slide 21 index 20）
    keys = [
        "关键技术一\n联邦本地子图与属性：在不出域下构建结构图与 kNN 属性语义超图。",
        "关键技术二\n跨视图对比与 InfoNCE：对齐属性视图与结构视图表示，温度系数与梯度分析见 §3.4。",
        "关键技术三\n全局语义原型：本地原型加噪上传、服务器聚合、回传正则抑制漂移。",
        "关键技术四\n语义感知压缩：元路径通道 Top-K + 残差量化与误差补偿。",
    ]

    # slide 22 难点（3 段 L2）
    diffs = [
        "难点一：零交互或极稀疏下结构视图退化；需可靠属性视图与对比对齐防噪声放大。",
        "难点二：Non-IID 下本地更新方向冲突；原型聚合与对比权重的平衡及稳定性。",
        "难点三：压缩误差与对比/元路径语义的耦合；需通道级语义选择与残差补偿防精度坍塌。",
    ]

    # slide 23 三列
    sol23 = [
        "难点对策一：双视图 + 全局原型，配合多任务损失与算法流程符号对齐（§3.7）。",
        "难点对策二：语义感知压缩管道与第四章实验 RQ 设计，控制变量说明 ACM 不设主表原因。",
        "难点对策三：系统侧分层实现与安全聚合、可选差分隐私，实测与理论通信一致。",
    ]

    # slide 24 四段 L_PARA
    sol24 = [
        "方法层：FedASCL 模块与联合优化；压缩层：选择器 + 量化 + 补偿；系统层：架构与模块实现。",
        "实验层：第三章总体/冷启动/Non-IID/消融；第四章压缩与帕累托；第五章功能与性能测试。",
        "评价层：HR/NDCG、通信字节与训练时间；表号与论文一致便于答辩翻阅。",
        "工程层：React/Flask/PyTorch/SQLite 等栈；联邦任务调度与监控面板。",
    ]

    # slide 28 四处「此处输入小标题」+ 正文 — 按整块替换含 L_MISC
    # slide 28 index 27 — 用子串替换简化：只替换 L_MISC 四次顺序
    misc28 = [
        "总体性能：tab:overall_performance；训练曲线 fig:training_progress。",
        "冷启动：tab:cold_start_res；基线集合与主表差异见 tab:chap3_exp_table_guide。",
        "Non-IID：tab:non_iid_data 与 fig:non_iid_robustness；极端 α 下衰减对比。",
        "消融：tab:ablation；验证属性视图与对比、原型项贡献。",
    ]

    # slide 29 成果应用四块
    app29 = [
        "添加文本\n算法成果：三数据集上 FedASCL 优于所列基线；Yelp 稀疏场景提升摘要所列幅度（以表为准）。",
        "添加文本\n压缩成果：约 20× 压缩比下精度可比；通信与挂钟时间量级与摘要一致。",
        "添加文本\n系统成果：推荐/设置/监控界面；联邦全流程与压缩上传实测。",
        "添加文本\n文档成果：硕士学位论文与可复现图表、表号体系。",
    ]

    # slide 31 课题结论 4 段 L1
    conc31 = [
        "提出 FedASCL：属性–结构双视图对比 + 全局语义原型，缓解冷启动与 Non-IID。",
        "提出语义感知压缩：元路径通道剪枝 + 残差量化 + 误差补偿，降低通信约一个数量级。",
        "实现联邦论文推荐系统：分层架构与端云闭环，实测与第四章结论互证。",
        "形成统一实验口径：主表、冷启动表、Non-IID 与表间指引，支撑可复核结论。",
    ]

    # slide 32 论文结论 + 01–03
    slide32_main = (
        "论文结论：本文围绕联邦异构图推荐的隐私与效果、通信与落地三方面，给出 FedASCL、语义感知压缩与系统实现，"
        "并在标准数据集与系统实测上完成验证；不足与展望见 §6.2。"
    )
    slide32_items = [
        "创新归纳：双视图 + 原型、通道级压缩 + 残差量化、工程闭环与隐私选项。",
        "实验归纳：RQ1–3 与消融覆盖主任务、冷启动与倾斜数据；压缩 RQ 覆盖精度–通信权衡。",
        "工程归纳：仿真引擎、冷启动路由、监控与通信压缩模块集成。",
    ]

    # slide 33 亮点 / 不足
    bright = (
        "亮点：问题—方法—实验—系统链条完整；表图与论文章节一一对应；压缩与系统实测相互印证。"
    )
    weak = (
        "不足：元路径选择与量化阈值为启发式；极端倾斜下原型加权可改进；动态图与更强隐私待拓展。"
    )

    # slide 34 未来方向三段
    fut = [
        "动态图与持续学习：在双视图与原型框架上建模兴趣随时间漂移。",
        "隐私增强：与 TEE、安全聚合及更强 DP 组合，降低结构侧推断风险。",
        "跨域少样本：在特征不完全重叠的联邦场景做迁移与对齐扩展。",
    ]

    # slide 35–36 参考文献占位：替换为首条代表性文献标题式（其余答辩时口述或换页）
    refs = [
        "Wu et al., FedGNN: Federated User Representation Learning with Graph Neural Networks (AAAI 2021).",
        "Tan et al., FedProto: Federated Prototype Learning across Heterogeneous Clients (ACM MM 2022).",
        "Yan et al., FedHGNN: Federated Heterogeneous Graph Neural Network (KDD 2024).",
        "Wang et al., KGAT: Knowledge Graph Attention Network for Recommendation (KDD 2019).",
        "He et al., LightGCN: Simplifying and Powering Graph Convolution Network for Recommendation (SIGIR 2020).",
        "GDPR; PIPL; 以及论文 chapters/reference.bib 中联邦推荐、对比学习、压缩相关条目（答辩可展开）。",
    ]

    counters: dict[tuple[int, str], int] = defaultdict(int)

    def scrub(s: str) -> str:
        """去掉模板残留短语（不改变已写入的论文章句主体）。"""
        if not s:
            return s
        for z in (L_YOUR, L_YOUR1, L_YOUR2):
            s = s.replace(z, "")
        for z in (
            "添加标题内容一",
            "添加标题内容二",
            "添加标题内容三",
            "添加标题内容四",
            "添加标题内容五",
            "添加标题内容",
        ):
            s = s.replace(z, "")
        s = re.sub(r"[ \t]+\n", "\n", s)
        lines = [ln.strip() for ln in s.split("\n")]
        s = "\n".join(ln for ln in lines if ln)
        s = re.sub(r"\n{3,}", "\n\n", s)
        return s.strip()

    def take(seq: list[str], si: int, key: str) -> str:
        k = (si, key)
        i = counters[k]
        counters[k] += 1
        return seq[i] if i < len(seq) else seq[-1]

    def transform_shape_text(si: int, t: str) -> str:
        t = global_sub(t)
        t = fill_slide_4(t, si)
        t = fill_slide_5_plain(t, si)
        t = fill_slide_6_double(t, si)

        # 标题微调（模板编号重复）
        if si == 7 and t.strip() == "1.5 主要创新点":
            return "创新点展开（一）"
        if si == 8 and t.strip() == "1.6 主要创新点":
            return "创新点展开（二）"

        if t == L1:
            # 通用：按页分发（slide 12–15、25、27 等大量 L1）
            pools: dict[int, list[str]] = {
                11: [  # 2.2 四段
                    "问题形式化：联邦本地子图、属性特征与 BPR 排序目标（§3.1）。",
                    "总体框架：服务器全局模型 + 全局原型；客户端双图编码与对比（§3.2）。",
                    "模块分解：属性图构建、跨视图对比、原型对齐、元路径推荐与联合损失（§3.3–3.6）。",
                    "算法与复杂度：alg:fedasclrc 流程对应；时间/空间与边规模同量级（§3.7–3.8）。",
                ],
                12: [  # 2.3 大段 — 单形状多段，用一段概括
                    "研究方法总述：在第二章符号与联邦范式基础上，采用「建模—优化—验证」路线；"
                    "第三章给出可学习的目标函数与联邦轮次；第四、五章给出压缩与系统约束下的实现与评测。",
                ],
                13: [  # 2.4 三列「添加相关方法i概述」+ L1 — 按出现顺序替换整段
                    "相关方法 A：kNN 语义超图与超边矩阵 H；余弦相似度与同类型节点约束（§3.3）。",
                    "相关方法 B：InfoNCE 跨视图对齐；温度 τ 与梯度传播（§3.4）。",
                    "相关方法 C：本地原型噪声、服务器聚合、全局原型正则（§3.5）。",
                    "相关方法 D：元路径语义注意力 + BPR 与多任务系数（§3.6）。",
                ],
                24: [  # 3.4 验证 — 三段 L2
                    "验证一：主实验与消融支持第三章结论；关键表号与论文一致。",
                    "验证二：压缩实验支持约 20× 设定与通信、帕累托结论。",
                    "验证三：系统实测流量/耗时与第四章同量级，闭环可信。",
                ],
                26: [  # 4.1 可能空或短 — 若出现 L2
                    "研究成果概述：算法、压缩、系统三线均有定量或界面证据；详见摘要与第六§6.1。",
                ],
                30: [  # slide 31 四段 plain L1
                    conc31[0],
                    conc31[1],
                    conc31[2],
                    conc31[3],
                ],
            }
            if si in pools:
                return take(pools[si], si, "l1")

        if t == L2 and si == 21:
            return take(diffs, si, "l2")
        if t == L2 and si == 24:
            return take(
                [
                    "验证（一）：对照第三章 tab:overall_performance、fig:training_progress 与消融表。",
                    "验证（二）：对照第四章压缩主表、通信成本与帕累托曲线。",
                ],
                si,
                "l2v",
            )
        if t == L2 and si == 25:
            return "验证（三）：对照第五章系统架构、界面与实测日志，与第四章通信量级一致。"

        if si == 10 and t.strip() in ("01 添加标题内容", "02 添加标题内容", "03 添加标题内容", "04 添加标题内容"):
            m = {
                "01 添加标题内容": "01 总体路线",
                "02 添加标题内容": "02 数据与符号",
                "03 添加标题内容": "03 训练与评测",
                "04 添加标题内容": "04 压缩与系统",
            }
            return m[t.strip()]

        if si == 13 and t.strip() == "添加相关方法i概述":
            return take(["A 属性语义图", "B 跨视图对比", "C 原型对齐", "D 元路径主任务"], si, "ov")

        if si == 14 and L_YOUR2 in t:
            return (
                "2.5 方法小结：在联邦与异构语义约束下，联合优化 BPR、InfoNCE 与原型项；"
                "上传阶段叠加通道剪枝与残差量化；系统侧完成全流程调度与监控。"
            )

        if si == 18 and L_YOUR2 in t:
            return (
                "2.3 研究方法与路线（续）：以公开数据集与统一划分为基础，"
                "先完成第三章全精度主实验，再接入第四章压缩与第五章系统实测，保证结论可对照。"
            )

        if si == 27 and t.strip() == "此处输入小标题":
            subs = ["总体性能（RQ1）", "冷启动（RQ2）", "Non-IID（RQ3）", "消融实验"]
            return take(subs, si, "s28sub")

        if si == 27 and t.strip() == L_MISC:
            return take(misc28, si, "s28misc")

        if t == L_ADD_TITLE and si == 22:
            return take(["难点子项 A", "难点子项 B", "难点子项 C"], si, "adt")

        if t == L_COLOR and si == 22:
            return take(sol23, si, "color")

        if t == L_PARA and si == 23:
            return take(sol24, si, "para")

        if t == L_THESIS_CONC:
            return slide32_main

        if t == L_YOUR2 and si == 10:
            return take(slide8_bodies, si, "your2")
        if t == L_YOUR1 and si == 10:
            return take(slide8_bodies, si, "your1")

        # slide 11 index 10: 01 添加标题内容 + YOUR1
        if si == 10 and "添加标题内容" in t and L_YOUR1 in t:
            parts = [
                "01 总体路线\n问题：联邦异构图冷启动、Non-IID 与通信。方案：FedASCL + 语义感知压缩 + 系统。",
                "02 数据与符号\n三数据集、属性字段、客户端划分与基线；第二章给出图与联邦基础。",
                "03 训练与评测\n主表、冷启动、Non-IID、消融与曲线；统一表间说明。",
                "04 压缩与系统\n上传管道、实验 RQ；分层架构、仿真与界面验证。",
            ]
            return take(parts, si, "s11")

        if t == L_LOREM and si == 15:
            return take(exp_steps, si, "lorem")
        if L_LOREM_SHORT in t and "01.添加文本标题" in t and si == 16:
            # 整段替换
            return take(goals, si, "goal")

        if si == 16 and t.strip().startswith("02.添加文本标题"):
            return goals[1]
        if si == 16 and t.strip().startswith("03.添加文本标题"):
            return goals[2]
        if si == 16 and t.strip().startswith("04.添加文本标题"):
            return goals[3]

        if si == 17 and FEAS in t:
            return "可行性：算法复杂度与主流 GNN 可比；压缩在公开数据集上可复现；系统技术栈成熟，实验环境可部署。"

        if si == 17 and t.strip() == "添加文本标题":
            return take(["数据可行性", "算法可行性", "工程可行性", "评测可行性"], si, "feas_title")

        cn_key = ["一", "二", "三", "四"]
        for k in range(4):
            if f"关键技术{cn_key[k]}" in t and L_LOREM_SHORT in t:
                return keys[k]

        if si == 14 and "01\n添加标题" in t:
            return "01\n本地训练：FedASCL 多轮本地更新与损失计算。"
        if si == 14 and "02\n添加标题" in t:
            return "02\n压缩上传：元路径选择器 + 量化梯度与误差状态。"
        if si == 14 and "03\n添加标题" in t:
            return "03\n安全聚合：服务器加权更新全局参数与原型。"
        if si == 14 and "04\n添加标题" in t:
            return "04\n评测归档：曲线、表与通信日志对照论文图表。"

        if t == L_MISC:
            return take(misc28, si, "misc")

        if si == 27 and "添加文\n本标题" in t:
            return "论文实验与系统成果要点（对照第三～五章图表）"

        if "添加文\n本标题" in t:
            return take(["成果：总体性能", "成果：冷启动", "成果：Non-IID", "成果：消融"], si, "badwen")

        if "添加文本\n" in t and L_MISC2 in t:
            return take(app29, si, "app")

        if t == L_LONG_BAD:
            return take([bright, weak], si, "longbad")

        if t == L_CLICK:
            return take(fut, si, "click")

        if t == L_REF:
            return take(refs, si, "ref")

        # slide 7（index 6）组合内仅替换与模板完全一致的 L1 段落
        if si == 6 and t.strip() == L1:
            return take(innov, si, "innov")

        if si == 8 and t.startswith("0") and L_YOUR1 in t:
            # 01\n ... slide 9
            idx = {"01": 0, "02": 1, "03": 2, "04": 3}.get(t[:2], 0)
            return slide9_boxes[idx]

        if si == 7 and L_YOUR in t:
            return take(slide8_bodies, si, "s8")

        if si == 7 and t.startswith("添加标题内容"):
            return t

        if si == 31:
            if t == L1:
                return take(slide32_items, si, "s32")

        # slide 27 可能「添加本本标题」已 global_sub 成小节要点
        if si == 26 and L2 in t:
            return (
                "4.1 研究成果概述：第三章主实验与消融、第四章压缩与通信、第五章系统与实测，"
                "共同支撑摘要中的定量结论与工程可行性陈述。"
            )

        # slide 13 复杂块
        if si == 12 and "添加标题内容" in t and L_YOUR1 in t:
            return (
                "2.3 研究方法\n"
                "规范研究：文献综述（§1.2）+ 理论准备（第二章）；设计研究：FedASCL 与压缩（第三、四章）；"
                "实证研究：对比实验、消融与系统测试（第三～五章）。"
            )

        if si == 12 and t.strip() == "添加\n标题":
            return "方法主线"

        # slide 25 超长 L3
        if L3 in t:
            return (
                "验证分析：主实验复现论文表；压缩实验复现通信与帕累托；系统测试复现监控与日志；"
                "三者与 tab:chap3_exp_table_guide 口径一致。"
            )

        # slide 18 2.8
        if si == 17 and "添加文本标题" in t and FEAS not in t:
            return "2.8 可行性：数据与场景来自公开数据集；算法模块基于 PyG；系统基于 Flask/React 可部署。"

        # 致谢
        if si == 36 and "感谢母校" in t:
            return (
                "感谢西北大学提供的培养平台与科研条件；\n"
                "特别感谢导师吴昊老师在选题、方案与论文定稿中的悉心指导；\n"
                "感谢实验室同学与舍友的帮助与支持；\n"
                "感谢各位答辩专家批评指正。"
            )

        return t

    for si, slide in enumerate(prs.slides):

        def fn(sh):
            tf = sh.text_frame
            old = tf.text
            new = scrub(transform_shape_text(si, old))
            if new != old:
                tf.text = new

        walk_shapes(slide.shapes, fn)

    # 封面日期（原模板空文本框）
    s0 = prs.slides[0].shapes
    if len(s0) > 6 and hasattr(s0[6], "text_frame") and not s0[6].text_frame.text.strip():
        s0[6].text_frame.text = "二〇二六年三月"
    if len(s0) > 7 and hasattr(s0[7], "text_frame") and not s0[7].text_frame.text.strip():
        s0[7].text_frame.text = "硕士专业学位论文（电子信息）· 答辩稿"

    prs.save(str(DST))
    print("已写入:", DST)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
